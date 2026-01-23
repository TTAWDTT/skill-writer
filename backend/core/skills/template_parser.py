"""
模板文件解析器
支持解析 md, doc, docx, pdf, txt, pptx 等格式的模板文件
"""
import re
from typing import Dict, List, Any, Optional
from pathlib import Path


def parse_template_file(content: bytes, file_ext: str, filename: str) -> str:
    """
    解析模板文件内容

    Args:
        content: 文件二进制内容
        file_ext: 文件扩展名
        filename: 文件名

    Returns:
        解析后的文本内容
    """
    if file_ext in ['.md', '.txt']:
        # 直接解码文本文件
        try:
            return content.decode('utf-8')
        except UnicodeDecodeError:
            return content.decode('gbk', errors='ignore')

    elif file_ext == '.docx':
        return parse_docx(content)

    elif file_ext == '.doc':
        # .doc 格式较旧，尝试使用简单方法提取文本
        return parse_doc(content)

    elif file_ext == '.pptx':
        return parse_pptx(content)

    elif file_ext == '.pdf':
        return parse_pdf(content)

    else:
        raise ValueError(f"Unsupported file type: {file_ext}")


def parse_docx(content: bytes) -> str:
    """解析 DOCX 文件"""
    try:
        import zipfile
        import io
        from xml.etree import ElementTree as ET

        # DOCX 是 ZIP 格式
        with zipfile.ZipFile(io.BytesIO(content)) as zf:
            # 读取 document.xml
            with zf.open('word/document.xml') as doc:
                tree = ET.parse(doc)
                root = tree.getroot()

                # XML 命名空间
                namespaces = {
                    'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'
                }

                # 提取所有文本
                paragraphs = []
                for para in root.iter('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}p'):
                    texts = []
                    for text in para.iter('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t'):
                        if text.text:
                            texts.append(text.text)
                    if texts:
                        paragraphs.append(''.join(texts))

                return '\n\n'.join(paragraphs)

    except Exception as e:
        raise ValueError(f"Failed to parse DOCX file: {str(e)}")


def parse_doc(content: bytes) -> str:
    """解析 DOC 文件 (简单文本提取)"""
    try:
        # 尝试提取 DOC 文件中的可读文本
        # DOC 是二进制格式，这里使用简单的文本提取方法
        text_parts = []

        # 尝试解码为文本
        try:
            text = content.decode('utf-8', errors='ignore')
        except:
            text = content.decode('latin-1', errors='ignore')

        # 过滤出可打印字符
        import string
        printable = set(string.printable + '中文日本語한국어абвгдеёжзийклмнопрстуфхцчшщъыьэюя')

        # 提取连续的可读文本段
        current_text = []
        for char in text:
            if char in printable or '\u4e00' <= char <= '\u9fff':
                current_text.append(char)
            else:
                if len(current_text) > 10:  # 只保留较长的文本段
                    text_parts.append(''.join(current_text))
                current_text = []

        if current_text:
            text_parts.append(''.join(current_text))

        result = '\n'.join(text_parts)

        if not result.strip():
            raise ValueError("Could not extract text from DOC file. Please convert to DOCX format.")

        return result

    except Exception as e:
        raise ValueError(f"Failed to parse DOC file: {str(e)}. Please convert to DOCX format.")


def parse_pptx(content: bytes) -> str:
    """解析 PPTX 文件"""
    try:
        import io
        from pptx import Presentation

        presentation = Presentation(io.BytesIO(content))
        text_parts = []

        for slide in presentation.slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    for paragraph in shape.text_frame.paragraphs:
                        runs_text = "".join(run.text for run in paragraph.runs).strip()
                        text = runs_text or paragraph.text.strip()
                        if text:
                            text_parts.append(text)
                elif getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text:
                                text_parts.append(cell_text)

        return "\n".join(text_parts)

    except Exception as e:
        raise ValueError(f"Failed to parse PPTX file: {str(e)}")


def parse_pdf(content: bytes) -> str:
    """解析 PDF 文件"""
    try:
        import io

        # 优先使用 pypdf（更活跃的维护）
        try:
            from pypdf import PdfReader  # type: ignore

            reader = PdfReader(io.BytesIO(content))
            text_parts = []
            for page in reader.pages:
                text = page.extract_text() or ""
                if text.strip():
                    text_parts.append(text)
            return "\n\n".join(text_parts).strip()
        except Exception:
            pass

        # 兼容旧依赖 PyPDF2
        try:
            import PyPDF2  # type: ignore

            reader = PyPDF2.PdfReader(io.BytesIO(content))
            text_parts = []
            for page in reader.pages:
                text = page.extract_text() or ""
                if text.strip():
                    text_parts.append(text)
            return "\n\n".join(text_parts).strip()
        except Exception:
            pass

        # 最后退化：简单文本流提取（效果有限，扫描版/图片 PDF 可能为空）
        text = content.decode("latin-1", errors="ignore")
        text_parts = re.findall(r"\(([^)]+)\)", text)
        return "\n".join(text_parts).strip()

    except Exception as e:
        raise ValueError(f"Failed to parse PDF file: {str(e)}")


def extract_sections_from_text(text: str) -> List[Dict[str, Any]]:
    """
    从文本中提取章节结构

    Args:
        text: 模板文本内容

    Returns:
        章节列表
    """
    sections = []

    # 匹配常见的标题格式
    # 1. Markdown 格式: # 标题, ## 标题
    # 2. 数字编号: 1. 标题, 1.1 标题
    # 3. 中文编号: 一、标题, （一）标题
    # 4. 全角数字: １．标题

    lines = text.split('\n')
    current_section = None

    for line in lines:
        line = line.strip()
        if not line:
            continue

        # Markdown 标题
        md_match = re.match(r'^(#{1,6})\s+(.+)$', line)
        if md_match:
            level = len(md_match.group(1))
            title = md_match.group(2).strip()
            section_id = re.sub(r'[^a-z0-9]', '-', title.lower())
            sections.append({
                'id': section_id,
                'title': title,
                'level': level,
                'type': 'required',
                'description': '',
                'children': []
            })
            continue

        # 数字编号标题 (1. 标题, 1.1 标题, 1.1.1 标题)
        num_match = re.match(r'^(\d+(?:\.\d+)*)[.、．]\s*(.+)$', line)
        if num_match:
            num_parts = num_match.group(1).split('.')
            level = len(num_parts)
            title = num_match.group(2).strip()
            section_id = f"section-{num_match.group(1).replace('.', '-')}"
            sections.append({
                'id': section_id,
                'title': title,
                'level': level,
                'type': 'required',
                'description': '',
                'children': []
            })
            continue

        # 中文编号标题 (一、标题, 二、标题)
        cn_num_match = re.match(r'^([一二三四五六七八九十]+)[、．.]\s*(.+)$', line)
        if cn_num_match:
            cn_nums = '一二三四五六七八九十'
            num_str = cn_num_match.group(1)
            level = 1
            title = cn_num_match.group(2).strip()
            section_id = f"section-cn-{cn_nums.index(num_str[0]) + 1}"
            sections.append({
                'id': section_id,
                'title': title,
                'level': level,
                'type': 'required',
                'description': '',
                'children': []
            })
            continue

        # 括号编号 ((一)标题, (1)标题)
        bracket_match = re.match(r'^[（(]([一二三四五六七八九十\d]+)[)）]\s*(.+)$', line)
        if bracket_match:
            level = 2
            title = bracket_match.group(2).strip()
            section_id = f"section-bracket-{bracket_match.group(1)}"
            sections.append({
                'id': section_id,
                'title': title,
                'level': level,
                'type': 'required',
                'description': '',
                'children': []
            })
            continue

    # 如果没有检测到章节，创建默认章节
    if not sections:
        sections = [
            {'id': 'introduction', 'title': '引言/概述', 'level': 1, 'type': 'required', 'description': '', 'children': []},
            {'id': 'main-content', 'title': '主体内容', 'level': 1, 'type': 'required', 'description': '', 'children': []},
            {'id': 'conclusion', 'title': '结论/总结', 'level': 1, 'type': 'required', 'description': '', 'children': []},
        ]

    return sections


def extract_fields_from_text(text: str) -> List[Dict[str, Any]]:
    """
    从模板文本中提取需要填写的字段

    Args:
        text: 模板文本内容

    Returns:
        字段列表
    """
    fields = []

    # 匹配常见的占位符格式
    # 1. [字段名], 【字段名】
    # 2. {{字段名}}, {字段名}
    # 3. <字段名>
    # 4. ____（下划线表示需要填写）

    # 中括号占位符
    bracket_fields = re.findall(r'[【\[]([\u4e00-\u9fff\w\s]+)[】\]]', text)

    # 花括号占位符
    brace_fields = re.findall(r'\{\{?([\u4e00-\u9fff\w\s]+)\}?\}', text)

    # 合并并去重
    all_field_names = list(set(bracket_fields + brace_fields))

    for i, name in enumerate(all_field_names[:15]):  # 最多15个字段
        field_id = re.sub(r'[^a-z0-9]', '_', name.lower())
        fields.append({
            'id': field_id,
            'name': name,
            'description': f'请输入{name}',
            'type': 'text' if len(name) < 10 else 'textarea',
            'required': True,
            'placeholder': f'请输入{name}...'
        })

    # 如果没有检测到字段，创建默认字段
    if not fields:
        fields = [
            {'id': 'title', 'name': '标题/名称', 'description': '文档的标题或名称', 'type': 'text', 'required': True, 'placeholder': '请输入标题...'},
            {'id': 'author', 'name': '作者/单位', 'description': '作者姓名或单位名称', 'type': 'text', 'required': False, 'placeholder': '请输入作者或单位...'},
            {'id': 'content_overview', 'name': '内容概述', 'description': '简要描述文档的主要内容', 'type': 'textarea', 'required': True, 'placeholder': '请简要描述主要内容...'},
        ]

    return fields


def generate_skill_from_template(
    template_content: str,
    skill_id: str,
    name: str,
    description: str,
    category: str,
    tags: List[str]
) -> Dict[str, Any]:
    """
    从模板内容生成 Skill 配置

    Args:
        template_content: 模板文本内容
        skill_id: Skill ID
        name: Skill 名称
        description: 描述
        category: 分类
        tags: 标签列表

    Returns:
        Skill 配置字典
    """
    # 提取章节结构
    sections = extract_sections_from_text(template_content)

    # 提取字段
    fields = extract_fields_from_text(template_content)

    # 生成写作指南
    guidelines = f"""# {name} 写作指南

## 文档概述
{description or f'本 Skill 用于帮助撰写 {name} 类型的文档。'}

## 章节结构
文档包含以下主要章节：
{chr(10).join(f'- {s["title"]}' for s in sections if s['level'] == 1)}

## 写作要求
1. 遵循模板的基本结构
2. 确保内容完整、逻辑清晰
3. 使用专业、规范的语言

## 注意事项
- 请根据实际情况调整内容
- 保持格式的一致性
"""

    # 生成系统提示词
    system_prompt = f"""你是一位专业的文档撰写专家，擅长撰写 {name} 类型的文档。

你的任务是帮助用户完成高质量的 {name} 文档撰写。

## 文档结构
{chr(10).join(f'{i+1}. {s["title"]}' for i, s in enumerate(sections) if s['level'] == 1)}

## 写作原则
1. 内容要准确、专业
2. 语言要规范、流畅
3. 结构要清晰、完整
4. 格式要统一、美观

请根据用户提供的信息，逐步完成文档的各个部分。
"""

    # 生成章节提示词模板
    section_prompt = """请撰写"{{ section_title }}"部分。

## 项目信息
{% for key, value in requirements.items() %}
- {{ key }}: {{ value }}
{% endfor %}

## 章节要求
{{ section_description }}

## 写作指导
{{ section_writing_guide }}

## 字数要求
{{ section_word_limit }}

请直接输出该章节的内容，使用 Markdown 格式。
"""

    return {
        'sections': sections,
        'fields': fields,
        'guidelines': guidelines,
        'system_prompt': system_prompt,
        'section_prompt': section_prompt,
        'instructions': f"""本 Skill 帮助您撰写 {name} 文档。

基于您上传的模板，系统已自动识别文档结构和需要填写的字段。
在对话过程中，请根据提示提供必要的信息，AI 将帮助您完成文档撰写。

## 文档结构
{chr(10).join(f'- {s["title"]}' for s in sections if s['level'] == 1)}
"""
    }
